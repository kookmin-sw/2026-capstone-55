"""
장풍소우 매장 사진 3장 - 원본 비율 유지하면서 깔끔하게 배치
각 사진의 원본 비율을 계산해서 잘리지 않도록 배치
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from PIL import Image
import os

# 사진 경로
photo_dir = os.path.expanduser("~/OneDrive/바탕 화면/장풍소우")
image_files = [
    os.path.join(photo_dir, "장풍소우1.jpg"),
    os.path.join(photo_dir, "장풍소우2.webp"),
    os.path.join(photo_dir, "장풍소우3.jpg"),
]

# webp 변환 및 원본 크기 확인
converted_files = []
image_sizes = []
for f in image_files:
    img = Image.open(f)
    image_sizes.append((img.width, img.height))
    print(f"{os.path.basename(f)}: {img.width}x{img.height} (비율 {img.width/img.height:.2f})")
    if f.lower().endswith('.webp'):
        new_path = f.replace('.webp', '_converted.png')
        img.save(new_path, 'PNG')
        converted_files.append(new_path)
    else:
        converted_files.append(f)

# PPT 생성 (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 배경 흰색
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# === 비율 유지 배치 함수 ===
def fit_image_in_box(img_w, img_h, box_w, box_h):
    """이미지를 박스 안에 비율 유지하면서 최대 크기로 맞춤"""
    img_ratio = img_w / img_h
    box_ratio = box_w / box_h
    
    if img_ratio > box_ratio:
        # 이미지가 더 넓음 -> 너비에 맞춤
        final_w = box_w
        final_h = box_w / img_ratio
    else:
        # 이미지가 더 높음 -> 높이에 맞춤
        final_h = box_h
        final_w = box_h * img_ratio
    
    # 박스 내 중앙 정렬 오프셋
    offset_x = (box_w - final_w) / 2
    offset_y = (box_h - final_h) / 2
    
    return final_w, final_h, offset_x, offset_y

# === 레이아웃 설정 ===
margin_x = Inches(0.4)
margin_y = Inches(0.35)
gap = Inches(0.15)

slide_w = prs.slide_width
slide_h = prs.slide_height
usable_w = slide_w - 2 * margin_x
usable_h = slide_h - 2 * margin_y

# 상단: 높이 57%, 하단: 나머지
top_box_w = usable_w
top_box_h = usable_h * 0.57

bottom_box_w = (usable_w - gap) / 2
bottom_box_h = usable_h * 0.43 - gap

# === 상단 사진 (비율 유지) ===
w1, h1 = image_sizes[0]
fw, fh, ox, oy = fit_image_in_box(w1, h1, top_box_w, top_box_h)
slide.shapes.add_picture(
    converted_files[0],
    Emu(int(margin_x + ox)), Emu(int(margin_y + oy)),
    Emu(int(fw)), Emu(int(fh))
)

# === 하단 왼쪽 사진 (비율 유지) ===
bottom_top = margin_y + top_box_h + gap
w2, h2 = image_sizes[1]
fw2, fh2, ox2, oy2 = fit_image_in_box(w2, h2, bottom_box_w, bottom_box_h)
slide.shapes.add_picture(
    converted_files[1],
    Emu(int(margin_x + ox2)), Emu(int(bottom_top + oy2)),
    Emu(int(fw2)), Emu(int(fh2))
)

# === 하단 오른쪽 사진 (비율 유지) ===
w3, h3 = image_sizes[2]
fw3, fh3, ox3, oy3 = fit_image_in_box(w3, h3, bottom_box_w, bottom_box_h)
right_start = margin_x + bottom_box_w + gap
slide.shapes.add_picture(
    converted_files[2],
    Emu(int(right_start + ox3)), Emu(int(bottom_top + oy3)),
    Emu(int(fw3)), Emu(int(fh3))
)

# === 저장 ===
output_path = os.path.join(os.path.expanduser("~/OneDrive/바탕 화면"), "장풍소우_매장소개.pptx")
prs.save(output_path)

print(f"\n✅ PPT 생성 완료! (비율 유지)")
print(f"📁 저장 위치: {output_path}")
