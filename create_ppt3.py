"""
장풍소우 매장소개 3 - 프리미엄 세련된 디자인
- 검은 배경
- 로고 중앙 상단
- 사진 3장 세련되게 배치
- 고급스러운 폰트 스타일
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import os

# 경로
photo_dir = os.path.expanduser("~/OneDrive/바탕 화면/장풍소우")
image_files = [
    os.path.join(photo_dir, "장풍소우1.jpg"),
    os.path.join(photo_dir, "장풍소우2.webp"),
    os.path.join(photo_dir, "장풍소우3.jpg"),
]
logo_path = os.path.join(photo_dir, "장풍소우로고.png")

# webp 변환
converted_files = []
for f in image_files:
    img = Image.open(f)
    print(f"{os.path.basename(f)}: {img.width}x{img.height}")
    if f.lower().endswith('.webp'):
        new_path = f.replace('.webp', '_converted.png')
        img.save(new_path, 'PNG')
        converted_files.append(new_path)
    else:
        converted_files.append(f)

# 크롭 함수
def crop_to_fill(img_path, target_w, target_h, output_path):
    img = Image.open(img_path)
    img_w, img_h = img.size
    target_ratio = target_w / target_h
    img_ratio = img_w / img_h
    if img_ratio > target_ratio:
        new_w = int(img_h * target_ratio)
        new_h = img_h
        left = (img_w - new_w) // 2
        top = 0
    else:
        new_w = img_w
        new_h = int(img_w / target_ratio)
        left = 0
        top = (img_h - new_h) // 2
    cropped = img.crop((left, top, left + new_w, top + new_h))
    cropped.save(output_path, quality=95)
    return output_path

os.makedirs("temp_crops", exist_ok=True)

# PPT 생성
prs = Presentation()
slide_w = Inches(13.333)
slide_h = Inches(7.5)
prs.slide_width = slide_w
prs.slide_height = slide_h

slide = prs.slides.add_slide(prs.slide_layouts[6])

# 배경 완전 검정
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# === 레이아웃 구성 ===
# 상단: 로고 영역 (중앙)
# 중간: 텍스트 (캐치프레이즈)
# 하단: 사진 3장 나란히
# 맨 하단: 연락처

logo_area_top = Inches(0.3)
logo_h = Inches(1.6)

text_top = logo_area_top + logo_h + Inches(0.15)
text_h = Inches(0.7)

photo_top = text_top + text_h + Inches(0.15)
photo_gap = Inches(0.08)
bottom_bar_h = Inches(0.65)
photo_h = slide_h - photo_top - bottom_bar_h - Inches(0.15)

# === 로고 (중앙 배치) ===
logo_img = Image.open(logo_path)
logo_w_px, logo_h_px = logo_img.size
logo_ratio = logo_w_px / logo_h_px
logo_display_h = logo_h
logo_display_w = Emu(int(logo_display_h * logo_ratio))

# 중앙 정렬
logo_x = Emu(int((slide_w - logo_display_w) / 2))
logo_y = Emu(int(logo_area_top))
slide.shapes.add_picture(logo_path, logo_x, logo_y, logo_display_w, logo_display_h)

# === 캐치프레이즈 텍스트 (중앙) ===
txBox = slide.shapes.add_textbox(Emu(0), Emu(int(text_top)), slide_w, Emu(int(text_h)))
tf = txBox.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER

run1 = p.add_run()
run1.text = "눈으로 직접보고 고르는 정육식당"
run1.font.size = Pt(22)
run1.font.bold = False
run1.font.color.rgb = RGBColor(212, 175, 55)  # 금색
run1.font.name = "맑은 고딕"

run2 = p.add_run()
run2.text = "   |   "
run2.font.size = Pt(18)
run2.font.color.rgb = RGBColor(80, 80, 80)

run3 = p.add_run()
run3.text = "150석 완비"
run3.font.size = Pt(22)
run3.font.bold = False
run3.font.color.rgb = RGBColor(255, 255, 255)
run3.font.name = "맑은 고딕"

# === 사진 3장 (가로 나란히, 동일 크기) ===
num_photos = 3
total_gap = photo_gap * (num_photos - 1)
each_photo_w = (slide_w - Inches(0.6) - total_gap) / num_photos  # 좌우 여백 0.3씩
photo_start_x = Inches(0.3)

for i, conv_file in enumerate(converted_files):
    ext = ".png" if conv_file.endswith('.png') else ".jpg"
    crop_path = f"temp_crops/c{i+1}{ext}"
    cropped = crop_to_fill(conv_file, each_photo_w, photo_h, crop_path)
    
    x = Emu(int(photo_start_x + i * (each_photo_w + photo_gap)))
    slide.shapes.add_picture(cropped, x, Emu(int(photo_top)), Emu(int(each_photo_w)), Emu(int(photo_h)))

# === 하단 연락처 (세련된 스타일) ===
bottom_top = slide_h - bottom_bar_h - Inches(0.05)
txBox2 = slide.shapes.add_textbox(Emu(0), Emu(int(bottom_top)), slide_w, Emu(int(bottom_bar_h)))
tf2 = txBox2.text_frame
tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
p2 = tf2.paragraphs[0]
p2.alignment = PP_ALIGN.CENTER

run4 = p2.add_run()
run4.text = "TEL. 032-831-1112"
run4.font.size = Pt(16)
run4.font.bold = False
run4.font.color.rgb = RGBColor(212, 175, 55)
run4.font.name = "맑은 고딕"

run5 = p2.add_run()
run5.text = "    "
run5.font.size = Pt(16)

run6 = p2.add_run()
run6.text = "송도 컨벤시아 남문 건너편"
run6.font.size = Pt(14)
run6.font.color.rgb = RGBColor(170, 170, 170)
run6.font.name = "맑은 고딕"

# === 저장 ===
output_path = os.path.join(os.path.expanduser("~/OneDrive/바탕 화면"), "장풍소우_매장소개4.pptx")
prs.save(output_path)
print(f"\n✅ PPT 생성 완료!")
print(f"📁 저장 위치: {output_path}")

import shutil
shutil.rmtree("temp_crops", ignore_errors=True)
