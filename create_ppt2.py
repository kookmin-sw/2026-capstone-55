"""
장풍소우 매장소개 2 - 여백 없이 꽉 채우면서 비율 최대한 유지
사진을 슬라이드 전체에 빈틈없이 채움 (crop 방식)
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from PIL import Image
import os

# 사진 경로
photo_dir = os.path.expanduser("~/OneDrive/바탕 화면/장풍소우")
image_files = [
    os.path.join(photo_dir, "장풍소우1.jpg"),
    os.path.join(photo_dir, "장풍소우2.webp"),
    os.path.join(photo_dir, "장풍소우3.jpg"),
]

# webp 변환 및 원본 크기 확인
converted_files = []
image_sizes = []
for f in image_files:
    img = Image.open(f)
    image_sizes.append((img.width, img.height))
    print(f"{os.path.basename(f)}: {img.width}x{img.height} (비율 {img.width/img.height:.2f})")
    if f.lower().endswith('.webp'):
        new_path = f.replace('.webp', '_converted.png')
        img.save(new_path, 'PNG')
        converted_files.append(new_path)
    else:
        converted_files.append(f)

# 사진을 박스에 꽉 채우도록 crop (Pillow로 미리 잘라서 저장)
def crop_to_fill(img_path, target_w, target_h, output_path):
    """이미지를 target 비율에 맞게 중앙 크롭 후 저장"""
    img = Image.open(img_path)
    img_w, img_h = img.size
    target_ratio = target_w / target_h
    img_ratio = img_w / img_h
    
    if img_ratio > target_ratio:
        # 이미지가 더 넓음 -> 좌우 크롭
        new_w = int(img_h * target_ratio)
        new_h = img_h
        left = (img_w - new_w) // 2
        top = 0
    else:
        # 이미지가 더 높음 -> 상하 크롭
        new_w = img_w
        new_h = int(img_w / target_ratio)
        left = 0
        top = (img_h - new_h) // 2
    
    cropped = img.crop((left, top, left + new_w, top + new_h))
    cropped.save(output_path, quality=95)
    return output_path

# PPT 생성 (16:9)
prs = Presentation()
slide_w = Inches(13.333)
slide_h = Inches(7.5)
prs.slide_width = slide_w
prs.slide_height = slide_h

blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)

# 배경 검정 (혹시 미세한 틈이 보일 경우 대비)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# === 레이아웃: 여백 0, 간격 최소 ===
gap = Inches(0.05)  # 사진 사이 아주 얇은 구분선

# 상단: 슬라이드 높이의 58%
top_h = slide_h * 0.58
top_w = slide_w

# 하단: 나머지 높이, 좌우 반반
bottom_h = slide_h - top_h - gap
bottom_w = (slide_w - gap) / 2

# === 사진 크롭 후 배치 ===
os.makedirs("temp_crops", exist_ok=True)

# 상단 사진 크롭
crop1 = crop_to_fill(converted_files[0], top_w, top_h, "temp_crops/crop1.jpg")
slide.shapes.add_picture(crop1, Emu(0), Emu(0), Emu(int(top_w)), Emu(int(top_h)))

# 하단 왼쪽
bottom_top = top_h + gap
crop2 = crop_to_fill(converted_files[1], bottom_w, bottom_h, "temp_crops/crop2.png")
slide.shapes.add_picture(crop2, Emu(0), Emu(int(bottom_top)), Emu(int(bottom_w)), Emu(int(bottom_h)))

# 하단 오른쪽
right_x = bottom_w + gap
crop3 = crop_to_fill(converted_files[2], bottom_w, bottom_h, "temp_crops/crop3.jpg")
slide.shapes.add_picture(crop3, Emu(int(right_x)), Emu(int(bottom_top)), Emu(int(bottom_w)), Emu(int(bottom_h)))

# === 저장 ===
output_path = os.path.join(os.path.expanduser("~/OneDrive/바탕 화면"), "장풍소우_매장소개2.pptx")
prs.save(output_path)

print(f"\n✅ PPT 생성 완료! (여백 없이 꽉 채움)")
print(f"📁 저장 위치: {output_path}")

# 임시 파일 정리
import shutil
shutil.rmtree("temp_crops", ignore_errors=True)
